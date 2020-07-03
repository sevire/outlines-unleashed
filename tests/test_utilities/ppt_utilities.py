from pptx import Presentation


def get_slide_data(file_name):
    """
    Creates a data structure from a .ppt file which allows tests to check whether a slide deck has been correctly
    created.  Assumes a fairly simple structure (which is ok at the time of writing)

    Extracts:
    - One record for each slide
    - Title for each slide
    - Sub-title for each slide
    - Bullets for each slide

    :return:
    """
    slides = Presentation(file_name).slides

    for slide in slides:
        shapes = slide.shapes
        for shape in shapes:
            text_frame = shape.text_frame
            paragraphs = text_frame.paragraphs
            for paragraph in paragraphs:
                level = paragraph.level
                runs = paragraph.runs
                text = ""
                for run in runs:
                    text += run.text
                yield (level, text)
