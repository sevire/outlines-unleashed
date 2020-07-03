"""
Functionality to transform a data object extracted from an outline to powerpoint files (or specific slides)
"""
from pptx import Presentation

powerpoint_slide_structure = {
    'deck_title': {
        'layout_index': 0,
        'placeholders': {
            'deck_title': 0,
            'deck_subtitle': 1
        }
    },
    'section_header': {
        'layout_index': 2,
        'placeholders': {
            'section_header': 0,
            'section_header_text': 1
        }
    },
    'content': {
        'layout_index': 1,
        'placeholders': {
            'content_title': 0,
            'content_bullets': 1
        }
    }
}


class PowerPointGenerator:
    @staticmethod
    def _print_placeholders_in_slide(slide):
        """
        Mainly used for debugging.  Allows us to see what is in a slide which may help track down
        errors in trying to add new elements (e.g. Missing layouts or placeholders)

        :param slide:
        :return:
        """
        for shape in slide.placeholders:
            print('%d %s' % (shape.placeholder_format.idx, shape.name))

    @staticmethod
    def _get_layout_index(slide_type):
        """
        PowerPoint represents the slide master layouts with an index.  This method
        will work out the right layout index for the given slide type represented
        within the extracted_data_node_table.

        :param slide_type:
        :return: the index to use within python-pptx to add the slide of the
                 given type
        """
        if slide_type in powerpoint_slide_structure:
            return powerpoint_slide_structure[slide_type]['layout_index']
        else:
            return None

    @staticmethod
    def _get_placeholder_index(slide_type, placeholder_name):
        if slide_type in powerpoint_slide_structure:
            slide_data = powerpoint_slide_structure[slide_type]
            if placeholder_name in slide_data['placeholders']:
                return slide_data['placeholders'][placeholder_name]

    @staticmethod
    def _add_title_slide(presentation, deck_title="(un-specified)", deck_sub_title="(un-specified)"):
        """
        Add the overall title slide for a slide deck

        There are potentially two fields to add; Title and Sub-title.

        :param deck_title:
        :param deck_sub_title:
        :return:
        """
        title_slide_layout_index = PowerPointGenerator._get_layout_index('deck_title')
        title_slide_layout = presentation.slide_layouts[title_slide_layout_index]
        slide = presentation.slides.add_slide(title_slide_layout)
        # PowerPointGenerator._print_placeholders_in_slide(slide)

        if deck_title is not None:
            title_placeholder_key = PowerPointGenerator._get_placeholder_index('deck_title', 'deck_title')
            title_placeholder = slide.placeholders[title_slide_layout_index]
            title_placeholder.text = deck_title

        if deck_sub_title is not None:
            subtitle_placeholder_key = PowerPointGenerator._get_placeholder_index('deck_title', 'deck_subtitle')
            subtitle_placeholder = slide.placeholders[subtitle_placeholder_key]
            subtitle_placeholder.text = deck_sub_title

    @staticmethod
    def _add_section_slide(presentation, section_title_text="(un-specified)", section_sub_title_text="(un-specified)"):
        """
        :param presentation:
        :param section_title_text:
        :param section_sub_title_text:
        :return:
        """
        section_title_slide_layout_index = PowerPointGenerator._get_layout_index('section_header')
        section_title_slide_layout = presentation.slide_layouts[section_title_slide_layout_index]
        slide = presentation.slides.add_slide(section_title_slide_layout)
        # PowerPointGenerator._print_placeholders_in_slide(slide)

        if section_title_text is not None:
            section_title_placeholder_key = PowerPointGenerator._get_placeholder_index('section_header', 'section_header')
            section_title = slide.shapes.placeholders[section_title_placeholder_key]
            section_title.text = section_title_text

        if section_sub_title_text is not None:
            section_subtitle_placeholder_key = PowerPointGenerator._get_placeholder_index('section_header', 'section_header_text')
            section_subtitle = slide.shapes.placeholders[section_subtitle_placeholder_key]
            section_subtitle.text = section_sub_title_text

    @staticmethod
    def _add_content_slide(presentation, title_text="(un-specified"):
        bullet_slide_layout_index = PowerPointGenerator._get_layout_index('content')
        bullet_slide_layout = presentation.slide_layouts[bullet_slide_layout_index]

        slide = presentation.slides.add_slide(bullet_slide_layout)
        # PowerPointGenerator._print_placeholders_in_slide(slide)

        shapes = slide.shapes

        title_placeholder_key = PowerPointGenerator._get_placeholder_index('content', 'content_title')
        title_shape = shapes.placeholders[title_placeholder_key]
        title_shape.text = title_text

        body_shape_key = PowerPointGenerator._get_placeholder_index('content', 'content_bullets')
        body_shape = shapes.placeholders[body_shape_key]
        text_frame = body_shape.text_frame

        return text_frame

    @staticmethod
    def _add_slide_bullet(text_frame, bullet_text, bullet_level, first_bullet=False):
        """
        Adds a slide of bullets to the presentation.

        :param text_frame:
        :param bullet_text:
        :param bullet_level:
        :param first_bullet: If the first bullet is being added, we don't need to add a paragraph as one would have
                             been created when the slide was created.
        :return:
        """

        if first_bullet is True:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.text = bullet_text
        paragraph.level = bullet_level

    def create_power_point_skeleton(self, data_node, template_ppt_path, output_ppt_path):
        """
        Takes an un-transformed outline (for now) and interprets it in order to create a PowerPoint slide deck.

        The outline is processed as follows:

        - Root Node:
          - Tag: Main Title for deck.
          - Text: Sub-title for deck.
        - Level 1 Node:
          - Tag: Section Title
          - Text: Section Sub-Title
        - Level 2 Node
          - Text: Slide Title
        - Level 2+n Node (n=1...)
          - Text: Bullet Level n

        The idea is to be able to create a skeleton presentation quickly.  The use case is professionals
        who use PPT as a status reporting tool (for example) who create a lot of presentations which aren't
        sophisticated but just structured notes into a slide form.

        Can also be used as the first iteration of a more sophisticated presentation.
        :return:
        """
        prs = Presentation(template_ppt_path)

        section = ""
        slide = ""
        text_frame = None  # Should be set before referenced as first row will always invoke a new slide

        first_bullet = None
        nodes = list(data_node.iter_unleashed_nodes())
        for node_record in nodes:
            node = node_record.node()
            if node_record.depth == 0:
                deck_title = node.text_tag
                deck_sub_title = node.text
                self._add_title_slide(prs, deck_title=deck_title, deck_sub_title=deck_sub_title)
            elif node_record.depth == 1:
                section_title = node.text_tag
                section_sub_title = node.text
                self._add_section_slide(prs, section_title_text=section_title, section_sub_title_text=section_sub_title)
            elif node_record.depth == 2:
                # Node is slide title, so we need to add a content slide and then process all subsequent bullet
                # nodes.
                title = node.text
                text_frame = self._add_content_slide(prs, title)
                first_bullet = True
            else:
                # Depth is more than two so this is a bullet to add to the existing text_frame

                bullet_depth = node_record.depth - 3
                bullet_text = node.text

                self._add_slide_bullet(text_frame, bullet_text, bullet_depth, first_bullet)
                first_bullet = False

        prs.save(output_ppt_path)
