from pptx import Presentation
from output_generators.ppt_output_generator_generic import PptOutputGeneratorGeneric


class PptOutputGeneratorSimple(PptOutputGeneratorGeneric):
    @staticmethod
    def generate_ppt(driver_table, output_path, template_file=None):
        # Initialisation
        prs = Presentation(template_file)

        # Parse each record and process accordingly
        pres_title, pres_sub_title, section_title, section_sub_title, slide_title, slide_sub_title, bullet = (None, None, None, None, None, None, None)
        text_frame = None
        first_bullet = None

        for record in driver_table:
            new_pres_title = record['slide_deck_title']
            new_pres_sub_title = record['slide_deck_sub_title']
            new_section_title = record['section_title']
            new_section_sub_title = record['section_sub_title']
            new_slide_title = record['slide_title']
            new_slide_sub_title = record['slide_sub_title']

            number_of_bullet_levels = 0
            bullet_values = []
            # Look for all the bullet fields by name (we know there won't be anywhere near 100 so using a loop
            for bullet_number in range(100):
                bullet_field_name = 'bullet-' + f"{bullet_number+1:02d}"
                if bullet_field_name not in record:
                    number_of_bullet_levels = bullet_number
                    break
                else:
                    bullet_value = record[bullet_field_name]
                    bullet_values.append(bullet_value)

            new_bullet_slide_required = False

            if new_pres_title != pres_title:
                # Probably the first slide so create deck title slide before processing rest of fields
                pres_title = new_pres_title
                new_bullet_slide_required = True
                PptOutputGeneratorSimple._add_title_slide(prs, deck_title=new_pres_title, deck_sub_title=new_pres_sub_title)

            if new_section_title != section_title:
                # This is a new section so create section title slide
                section_title = new_section_title
                new_bullet_slide_required = True
                PptOutputGeneratorSimple._add_section_slide(prs, new_section_title, new_section_sub_title)

            if new_slide_title != slide_title or new_bullet_slide_required is True:
                slide_title = new_slide_title
                first_bullet = True
                text_frame = PptOutputGeneratorSimple._add_content_slide(prs, new_slide_title)

            for bullet_number in range(number_of_bullet_levels):
                if bullet_values[bullet_number] is None:
                    # As soon as there is a None value stop adding bullet (it doesn't make sense to have a bullet
                    # which skips a level (although I believe PPT may allow it!)
                    break
                else:
                    level = bullet_number
                    PptOutputGeneratorSimple._add_slide_bullet(text_frame, bullet_values[bullet_number], level, first_bullet)
                    first_bullet = False
        prs.save(output_path)





